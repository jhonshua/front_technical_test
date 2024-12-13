from pptx import Presentation

# Abre la presentación existente
presentation = Presentation('trial.pptx')

# Recorre las diapositivas para encontrar y reemplazar el texto
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            if text_frame.text == "Hola":
                text_frame.text = "Gracias"

# Guarda los cambios en el archivo de presentación
presentation.save('trial.pptx')

print("Texto 'Hola' sustituido por 'Gracias' en 'trial.pptx'.")
